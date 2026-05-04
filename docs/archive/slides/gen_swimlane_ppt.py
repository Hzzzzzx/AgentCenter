#!/usr/bin/env python3
"""AgentCenter 信息流泳道图 PPT - 精确还原HTML 1200x700比例"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pptx.oxml.ns import qn

# 使用4:3标准比例 PPT (10x7.5)，更接近HTML的1200x700(1.714)
# 4:3 = 1.333，不对
# 让我们用 13.333 x 7.5 但保持原始1200x700的布局比例

# HTML: 1200 x 700 = 1.714
# PPT 16:9: 13.333 / 7.5 = 1.778 (接近!)

# 用统一缩放比例，保持1200x700的相对位置
SCALE = Inches(13.333) / 1200  # 统一缩放

def px(x): return x * SCALE  # px转inches

C = {
    'bg': RGBColor(0xFA,0xFA,0xFA), 'title': RGBColor(0x0A,0x0A,0x0A),
    'title_line': RGBColor(0x00,0x70,0xF3),
    'blue_bg': RGBColor(0xE3,0xF2,0xFD), 'green_bg': RGBColor(0xE8,0xF5,0xE9),
    'lane_border': RGBColor(0x90,0xA4,0xAE), 'label_border': RGBColor(0x19,0x76,0xD2),
    'lane_text': RGBColor(0x37,0x47,0x4F),
    'node_bg': RGBColor(0xFF,0xFF,0xFF), 'node_border': RGBColor(0x19,0x76,0xD2),
    'node_text': RGBColor(0x37,0x47,0x4F),
    'arrow': RGBColor(0x60,0x7D,0x8B),
    'agent_bg': RGBColor(0xEC,0xEF,0xF1), 'agent_border': RGBColor(0x60,0x7D,0x8B),
    'agent_node': RGBColor(0x00,0x70,0xF3), 'agent_node_border': RGBColor(0x0A,0x0A,0x0A),
    'agent_text': RGBColor(0xFF,0xFF,0xFF),
    'iline': RGBColor(0x78,0x90,0x9C), 'itext': RGBColor(0x78,0x90,0x9C),
    'pagenum': RGBColor(0x00,0x70,0xF3), 'dashed': RGBColor(0xB0,0xBE,0xC5),
}

# ===== 精确从HTML提取的数据 =====
lanes = [
    {
        "label": "需求系统", "bg": "blue_bg",
        "nodes": [(10,8,"提出需求"),(65,28,"需求评审"),(120,50,"需求分配"),(175,18,"需求确认")],
        "arrows": ["34,26 34,37 65,37","89,46 89,59 120,59","144,68 144,27 175,27"]
    },
    {
        "label": "开发环境", "bg": "green_bg",
        "nodes": [(420,8,"技术方案"),(475,50,"代码编写"),(530,8,"代码评审"),(585,50,"方案确认")],
        "arrows": ["444,26 444,50 475,50","499,59 499,26 530,26","554,26 554,50 585,50"]
    },
    {
        "label": "代码仓库", "bg": "blue_bg",
        "nodes": [(140,8,"提交代码"),(195,48,"发起MR"),(250,48,"代码审查"),(305,8,"合并代码")],
        "arrows": ["164,26 164,57 195,57","219,66 219,57 250,57","274,48 274,17 305,17"]
    },
    {
        "label": "CI/CD", "bg": "green_bg",
        "nodes": [(350,50,"触发构建"),(405,8,"执行测试"),(460,50,"生成报告"),(515,8,"镜像构建")],
        "arrows": ["374,68 374,17 405,17","429,26 429,59 460,59","484,68 484,17 515,17"]
    },
    {
        "label": "部署环境", "bg": "blue_bg",
        "nodes": [(20,8,"部署测试"),(75,52,"预发布"),(130,52,"发布生产"),(185,8,"健康检查")],
        "arrows": ["44,26 44,61 75,61","99,61 130,61","154,61 154,17 185,17"]
    },
    {
        "label": "运维管理", "bg": "green_bg",
        "nodes": [(320,52,"故障检测"),(375,8,"问题定位"),(430,8,"自动修复"),(485,52,"运维操作")],
        "arrows": ["344,61 344,17 375,17","399,17 430,17","454,17 454,61 485,61"]
    },
    {
        "label": "监控告警", "bg": "blue_bg",
        "nodes": [(480,12,"指标采集"),(535,48,"阈值检测"),(590,20,"触发告警"),(645,40,"告警确认")],
        "arrows": ["504,30 504,66 535,66","559,66 559,29 590,29","614,38 614,58 645,58"]
    },
    {
        "label": "通知协作", "bg": "green_bg",
        "nodes": [(200,12,"推送通知"),(255,32,"任务协作"),(310,52,"进度同步"),(365,16,"流程完成")],
        "arrows": ["224,30 224,41 255,41","279,50 279,70 310,70","334,70 334,25 365,25"]
    },
]

NO_MARKER_ARROWS = {"219,66 219,57 250,57", "99,61 130,61", "399,17 430,17"}

AGENT_NODES = ["信息感知", "智能分析", "任务调度", "执行反馈"]

INTERACTIONS = [
    (100, "调度指令", True),
    (170, "配置变更", True),
    (240, "任务分配", True),
    (350, "状态数据", False),
    (420, "告警信息", False),
    (490, "质量指标", False),
]


def add_polyline(slide, pts_str, ox, oy, marker=True):
    """绘制polyline折线箭头 - 只在最后一段的终点加箭头"""
    pts = []
    for p in pts_str.split():
        x, y = p.split(',')
        pts.append((float(x), float(y)))
    for i in range(len(pts)-1):
        x1,y1 = pts[i]
        x2,y2 = pts[i+1]
        c = slide.shapes.add_connector(1, px(ox+x1), px(oy+y1), px(ox+x2), px(oy+y2))
        c.line.color.rgb = C['arrow']
        c.line.width = Pt(1.5)
        # 只在最后一段加箭头，headEnd在连线的终点(end)加箭头
        if marker and i == len(pts)-2:
            ln = c.line._ln
            if ln is not None:
                he = etree.SubElement(ln, qn('a:headEnd'))
                he.set('type','triangle')
                he.set('w','med')
                he.set('len','med')


def add_node(slide, x, y, text):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(48), px(18))
    s.fill.solid(); s.fill.fore_color.rgb = C['node_bg']
    s.line.color.rgb = C['node_border']; s.line.width = Pt(1.5)
    s.adjustments[0] = 0.15
    tf = s.text_frame; tf.word_wrap = False; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(7); p.font.color.rgb = C['node_text']; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0); tf.paragraphs[0].space_after = Pt(0)
    for attr in ['margin_left','margin_right','margin_top','margin_bottom']:
        setattr(s.text_frame, attr, Pt(1))


def add_lane_label(slide, x, y, text, is_even):
    bg = 'green_bg' if is_even else 'blue_bg'
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(85), px(76))
    s.fill.solid(); s.fill.fore_color.rgb = C[bg]
    s.line.color.rgb = C['lane_border']; s.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(4), px(76))
    bar.fill.solid(); bar.fill.fore_color.rgb = C['label_border']; bar.line.fill.background()
    tf = s.text_frame; tf.word_wrap = False; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(10); p.font.color.rgb = C['lane_text']; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_agent_node(slide, x, y, text):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(200), px(55))
    s.fill.solid(); s.fill.fore_color.rgb = C['agent_node']
    s.line.color.rgb = C['agent_node_border']; s.line.width = Pt(2)
    s.adjustments[0] = 12/200
    tf = s.text_frame; tf.word_wrap = False; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(14); p.font.color.rgb = C['agent_text']; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_agent_arrow(slide, x, y):
    t = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, px(x), px(y), px(16), px(10))
    t.fill.solid(); t.fill.fore_color.rgb = C['iline']; t.line.fill.background()
    t.rotation = 180


def add_interaction(slide, x, y, label, left_arrow=True):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(60), Pt(2))
    line.fill.solid(); line.fill.fore_color.rgb = C['iline']; line.line.fill.background()
    ts = slide.shapes.add_textbox(px(x+30)-px(40), px(y-20), px(80), px(18))
    tf = ts.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = label
    p.font.size = Pt(9); p.font.color.rgb = C['itext']; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if left_arrow:
        a = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, px(x-6), px(y-4), px(8), px(10))
        a.rotation = 90
    else:
        a = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, px(x+58), px(y-4), px(8), px(10))
        a.rotation = -90
    a.fill.solid(); a.fill.fore_color.rgb = C['iline']; a.line.fill.background()


def add_dashed(slide, x, y, w):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), Pt(1))
    s.fill.solid(); s.fill.fore_color.rgb = C['dashed']; s.line.fill.background()


def create_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C['bg']

    # 标题
    ts = slide.shapes.add_textbox(px(25), px(12), px(500), px(30))
    tf = ts.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = "AgentCenter 信息流泳道图"
    p.font.size = Pt(24); p.font.color.rgb = C['title']; p.font.bold = True

    ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(25), px(42), px(100), px(3))
    ul.fill.solid(); ul.fill.fore_color.rgb = C['title_line']; ul.line.fill.background()

    # 泳道
    ll, lt = 25, 58  # lanes-area left, top
    cl = ll + 85  # content left
    for i, lane in enumerate(lanes):
        y = lt + i * 76
        add_lane_label(slide, ll, y, lane["label"], i%2==1)
        for nl, nt, nt_ in lane["nodes"]:
            add_node(slide, cl+nl, y+nt, nt_)
        for ar in lane["arrows"]:
            add_polyline(slide, ar, cl, y, ar not in NO_MARKER_ARROWS)

    # 虚线
    for i in range(1, 8):
        add_dashed(slide, cl, i*76-1, 695)

    # Agent面板
    pw, ph = 300, 608
    pl = 1200 - 15 - pw  # right=15
    pt = 58
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(pl), px(pt), px(pw), px(ph))
    panel.fill.solid(); panel.fill.fore_color.rgb = C['agent_bg']
    panel.line.color.rgb = C['agent_border']; panel.line.width = Pt(2)
    panel.line.dash_style = 4; panel.adjustments[0] = 12/300

    ts = slide.shapes.add_textbox(px(pl), px(pt), px(pw), px(45))
    tf = ts.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = "AgentCenter"
    p.font.size = Pt(18); p.font.color.rgb = C['title']; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    ts.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    sl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(pl), px(pt+45), px(pw), Pt(2))
    sl.fill.solid(); sl.fill.fore_color.rgb = C['agent_border']; sl.line.fill.background()
    sl.line.dash_style = 4

    # Agent节点
    ay = pt + 45
    for atext in AGENT_NODES:
        ay += 85
        ax = pl + (pw - 200) / 2
        add_agent_node(slide, ax, ay, atext)
        add_agent_arrow(slide, ax + (200-16)/2, ay + 55 + 2)
        ay += 55

    # 交互线
    for it, il, ia in INTERACTIONS:
        add_interaction(slide, 815, it, il, ia)

    # 页码
    pn = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(1180), px(654), px(34), px(34))
    pn.fill.solid(); pn.fill.fore_color.rgb = C['pagenum']; pn.line.fill.background()
    tf = pn.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = "5"
    p.font.size = Pt(13); p.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    pn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    create_slide(prs)
    output = "/Users/hzz/workspace/AgentCenter/docs/slides/swimlane.pptx"
    prs.save(output)
    print(f"✅ PPT已生成: {output}")


if __name__ == "__main__":
    main()
